"""
Extraction de données depuis des documents uploadés (PDF, DOCX, PPTX).
Utilise le LLM pour extraire les informations structurées.
"""

import io
import re
import json
from typing import Dict, Optional, List, Any
from dataclasses import dataclass, field
from enum import Enum
from loguru import logger

# PDF extraction
try:
    import fitz  # PyMuPDF
    HAS_PYMUPDF = True
except ImportError:
    HAS_PYMUPDF = False
    logger.warning("PyMuPDF not installed, PDF extraction disabled")

# DOCX extraction
try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    logger.warning("python-docx not installed, DOCX extraction disabled")

# PPTX extraction
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    logger.warning("python-pptx not installed, PPTX extraction disabled")


class DocumentType(Enum):
    """Types de documents supportés."""
    PDF = "pdf"
    DOCX = "docx"
    PPTX = "pptx"
    TXT = "txt"
    UNKNOWN = "unknown"


@dataclass
class ExtractedData:
    """Données extraites d'un document."""
    company_name: Optional[str] = None
    country: Optional[str] = None
    sector: Optional[str] = None
    subsector: Optional[str] = None
    business_description: Optional[str] = None
    employees: Optional[int] = None
    revenue: Optional[str] = None
    year_founded: Optional[int] = None
    target_market: Optional[str] = None
    geographic_scope: Optional[List[str]] = None
    
    # Données 2X Challenge
    women_ownership_pct: Optional[float] = None
    women_management_pct: Optional[float] = None
    women_employees_pct: Optional[float] = None
    benefits_women: Optional[bool] = None
    benefits_women_description: Optional[str] = None
    
    # Données produit/marché
    products_services: Optional[str] = None
    main_clients: Optional[str] = None
    competitive_advantage: Optional[str] = None
    
    # Métadonnées extraction
    raw_text: str = ""
    confidence: float = 0.0
    extraction_notes: List[str] = field(default_factory=list)
    source_filename: str = ""
    
    def to_dict(self) -> Dict:
        """Convertit en dictionnaire."""
        return {
            "company_name": self.company_name,
            "country": self.country,
            "sector": self.sector,
            "subsector": self.subsector,
            "business_description": self.business_description,
            "employees": self.employees,
            "revenue": self.revenue,
            "year_founded": self.year_founded,
            "target_market": self.target_market,
            "geographic_scope": self.geographic_scope,
            "women_ownership_pct": self.women_ownership_pct,
            "women_management_pct": self.women_management_pct,
            "women_employees_pct": self.women_employees_pct,
            "benefits_women": self.benefits_women,
            "benefits_women_description": self.benefits_women_description,
            "products_services": self.products_services,
            "main_clients": self.main_clients,
            "competitive_advantage": self.competitive_advantage,
            "confidence": self.confidence,
            "extraction_notes": self.extraction_notes,
        }
    
    def get_filled_fields_count(self) -> int:
        """Compte le nombre de champs remplis."""
        d = self.to_dict()
        return sum(1 for v in d.values() if v is not None and v != [] and v != "")
    
    def get_2x_data(self) -> Dict:
        """Retourne les données 2X Challenge."""
        return {
            "women_ownership_pct": self.women_ownership_pct or 0,
            "women_management_pct": self.women_management_pct or 0,
            "women_employees_pct": self.women_employees_pct or 0,
            "benefits_women": self.benefits_women or False,
        }


def detect_document_type(filename: str) -> DocumentType:
    """Détecte le type de document à partir de l'extension."""
    if not filename:
        return DocumentType.UNKNOWN
    
    ext = filename.lower().split('.')[-1]
    mapping = {
        'pdf': DocumentType.PDF,
        'docx': DocumentType.DOCX,
        'doc': DocumentType.DOCX,
        'pptx': DocumentType.PPTX,
        'ppt': DocumentType.PPTX,
        'txt': DocumentType.TXT,
        'md': DocumentType.TXT,
    }
    return mapping.get(ext, DocumentType.UNKNOWN)


def extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extrait le texte d'un PDF avec PyMuPDF."""
    if not HAS_PYMUPDF:
        return "[PDF extraction not available - install PyMuPDF]"
    
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        text_parts = []
        
        for page_num, page in enumerate(doc):
            page_text = page.get_text("text")
            if page_text.strip():
                text_parts.append(f"--- Page {page_num + 1} ---\n{page_text}")
        
        doc.close()
        return "\n\n".join(text_parts)
    
    except Exception as e:
        logger.error(f"Error extracting PDF text: {e}")
        return f"[Error extracting PDF: {str(e)}]"


def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extrait le texte d'un DOCX."""
    if not HAS_DOCX:
        return "[DOCX extraction not available - install python-docx]"
    
    try:
        doc = Document(io.BytesIO(file_bytes))
        text_parts = []
        
        # Extraire les paragraphes
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # Extraire les tableaux
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    table_text.append(row_text)
            if table_text:
                text_parts.append("\n[TABLE]\n" + "\n".join(table_text) + "\n[/TABLE]")
        
        return "\n\n".join(text_parts)
    
    except Exception as e:
        logger.error(f"Error extracting DOCX text: {e}")
        return f"[Error extracting DOCX: {str(e)}]"


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extrait le texte d'un PPTX."""
    if not HAS_PPTX:
        return "[PPTX extraction not available - install python-pptx]"
    
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"--- Slide {slide_num} ---"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
                
                # Extraire les tableaux des slides
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        if row_text.strip():
                            slide_text.append(row_text)
            
            if len(slide_text) > 1:  # Plus que juste le header
                text_parts.append("\n".join(slide_text))
        
        return "\n\n".join(text_parts)
    
    except Exception as e:
        logger.error(f"Error extracting PPTX text: {e}")
        return f"[Error extracting PPTX: {str(e)}]"


def extract_text(file_bytes: bytes, filename: str) -> str:
    """
    Extrait le texte d'un document selon son type.
    
    Args:
        file_bytes: Contenu du fichier en bytes
        filename: Nom du fichier (pour détecter le type)
    
    Returns:
        Texte extrait du document
    """
    doc_type = detect_document_type(filename)
    
    if doc_type == DocumentType.PDF:
        return extract_text_from_pdf(file_bytes)
    elif doc_type == DocumentType.DOCX:
        return extract_text_from_docx(file_bytes)
    elif doc_type == DocumentType.PPTX:
        return extract_text_from_pptx(file_bytes)
    elif doc_type == DocumentType.TXT:
        try:
            return file_bytes.decode('utf-8')
        except UnicodeDecodeError:
            return file_bytes.decode('latin-1', errors='ignore')
    else:
        logger.warning(f"Unknown document type for: {filename}")
        return f"[Unknown document type: {filename}]"


# =============================================================================
# Prompt d'extraction LLM
# =============================================================================

EXTRACTION_SYSTEM_PROMPT = """Tu es un assistant spécialisé dans l'extraction de données d'entreprise à partir de documents d'investissement (business plans, pitch decks, présentations).

Tu extrais les informations de manière structurée et précise. Si une information n'est pas trouvée dans le document, tu retournes null pour ce champ.

Tu travailles pour IPAE3, un fonds d'investissement à impact en Afrique subsaharienne."""


EXTRACTION_USER_PROMPT = """Analyse le texte suivant extrait d'un document d'entreprise et extrais les informations au format JSON.

INSTRUCTIONS STRICTES :
1. Extrais UNIQUEMENT les informations explicitement présentes dans le texte
2. Si une information n'est pas trouvée, utilise null
3. Pour les pourcentages, donne un nombre entier (ex: 35 pour 35%)
4. Pour le chiffre d'affaires, utilise ces tranches standardisées: "< 500K", "500K - 2M", "2M - 5M", "5M - 10M", "10M - 50M", "> 50M"
5. Pour le secteur, choisis UNIQUEMENT parmi: "Agribusiness", "Industrie & Manufacturing", "Services financiers", "Santé", "Éducation & Formation", "Énergie", "Tech & Digital", "Distribution & Retail", "Tourisme & Hôtellerie"
6. Pour target_market, choisis parmi: "B2C - Particuliers", "B2B - Entreprises", "B2B2C - Les deux", "B2G - Institutions"
7. Ajoute un score de confiance (0-100) basé sur la qualité et la quantité des informations trouvées

FORMAT DE SORTIE - Réponds UNIQUEMENT avec ce JSON, sans texte avant ou après :

```json
{{
  "company_name": "string ou null",
  "country": "string ou null (nom du pays)",
  "sector": "string ou null (parmi la liste)",
  "subsector": "string ou null",
  "business_description": "string ou null (2-3 phrases max décrivant l'activité)",
  "employees": "number ou null",
  "revenue": "string ou null (tranche standardisée)",
  "year_founded": "number ou null",
  "target_market": "string ou null (B2C/B2B/B2B2C/B2G)",
  "geographic_scope": ["liste des zones géographiques ou null"],
  "women_ownership_pct": "number ou null (0-100)",
  "women_management_pct": "number ou null (0-100)",
  "women_employees_pct": "number ou null (0-100)",
  "benefits_women": "boolean ou null",
  "benefits_women_description": "string ou null (si benefits_women=true, décrire comment)",
  "products_services": "string ou null (liste des produits/services principaux)",
  "main_clients": "string ou null",
  "competitive_advantage": "string ou null",
  "confidence": "number (0-100)",
  "extraction_notes": ["liste des points d'attention, incertitudes ou informations manquantes importantes"]
}}
```

TEXTE DU DOCUMENT À ANALYSER :
---
{document_text}
---

Réponds UNIQUEMENT avec le JSON, sans aucun texte supplémentaire."""


def parse_llm_extraction_response(llm_response: str) -> ExtractedData:
    """
    Parse la réponse du LLM et retourne un objet ExtractedData.
    
    Gère les cas où le JSON est dans un bloc de code ou mal formaté.
    """
    response = llm_response.strip()
    
    # Extraire le JSON s'il est dans un bloc de code
    json_match = re.search(r'```(?:json)?\s*(.*?)\s*```', response, re.DOTALL)
    if json_match:
        response = json_match.group(1)
    
    # Nettoyer les caractères problématiques
    response = response.strip()
    
    try:
        data = json.loads(response)
        
        return ExtractedData(
            company_name=data.get("company_name"),
            country=data.get("country"),
            sector=data.get("sector"),
            subsector=data.get("subsector"),
            business_description=data.get("business_description"),
            employees=data.get("employees"),
            revenue=data.get("revenue"),
            year_founded=data.get("year_founded"),
            target_market=data.get("target_market"),
            geographic_scope=data.get("geographic_scope"),
            women_ownership_pct=data.get("women_ownership_pct"),
            women_management_pct=data.get("women_management_pct"),
            women_employees_pct=data.get("women_employees_pct"),
            benefits_women=data.get("benefits_women"),
            benefits_women_description=data.get("benefits_women_description"),
            products_services=data.get("products_services"),
            main_clients=data.get("main_clients"),
            competitive_advantage=data.get("competitive_advantage"),
            confidence=data.get("confidence", 0),
            extraction_notes=data.get("extraction_notes", []),
        )
    
    except json.JSONDecodeError as e:
        logger.error(f"Failed to parse LLM response as JSON: {e}")
        logger.debug(f"Response was: {response[:500]}...")
        
        return ExtractedData(
            confidence=0,
            extraction_notes=[f"Échec du parsing JSON: {str(e)}"]
        )


def prepare_extraction_prompt(document_text: str, max_chars: int = 15000) -> str:
    """
    Prépare le prompt d'extraction avec troncature si nécessaire.
    
    Args:
        document_text: Texte complet du document
        max_chars: Nombre max de caractères à envoyer au LLM
    
    Returns:
        Prompt formaté
    """
    # Tronquer si trop long
    if len(document_text) > max_chars:
        # Garder le début et la fin (souvent les infos clés sont au début)
        half = max_chars // 2
        document_text = (
            document_text[:half] + 
            "\n\n[... DOCUMENT TRONQUÉ ...]\n\n" + 
            document_text[-half:]
        )
    
    return EXTRACTION_USER_PROMPT.format(document_text=document_text)


def get_extraction_system_prompt() -> str:
    """Retourne le system prompt pour l'extraction."""
    return EXTRACTION_SYSTEM_PROMPT
